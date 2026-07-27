# -*- coding: utf-8 -*-
"""
Gera apresentacao_defesa.pptx — defesa de mestrado (45 min, 16:9, pt-BR).
Filosofia: imagem narra, orador fala, texto é andaime.
Paleta (herdada do main.tex): navy 29,45,84 | azul 72,150,222 | vermelho 196,30,58 (SÓ resultados).
Rode: python prepare_assets.py && python build_deck.py
"""
from pathlib import Path
from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
A = HERE / 'assets'
DERIVED = A / '_derived'
DERIVED.mkdir(exist_ok=True)
OUT = HERE / 'apresentacao_defesa.pptx'

SW, SH = Inches(13.333), Inches(7.5)
NAVY = RGBColor(29, 45, 84)
BLUE = RGBColor(72, 150, 222)
RED = RGBColor(196, 30, 58)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(120, 120, 120)
DARK = RGBColor(40, 40, 40)
FONT = 'Calibri'

TITLE_H = Inches(0.95)


# ----------------------------------------------------------------------------
# infra
# ----------------------------------------------------------------------------
def _set_font(run, size, bold=False, color=DARK, italic=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def _txbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb


def _rect(slide, x, y, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def img_size(path):
    with Image.open(path) as im:
        return im.size


def fit(path, max_w, max_h):
    """(w,h) em EMU preservando proporção dentro da caixa."""
    w, h = img_size(path)
    scale = min(max_w / w, max_h / h)
    return int(w * scale), int(h * scale)


def add_pic_fit(slide, path, box_x, box_y, box_w, box_h):
    w, h = fit(path, box_w, box_h)
    x = box_x + int((box_w - w) / 2)
    y = box_y + int((box_h - h) / 2)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def darken_crop(src, out_name, target_ar=16 / 9, factor=0.35, navy_blend=0.45):
    """Center-crop para o aspect alvo, escurece e tinge de navy. Devolve o path."""
    out = DERIVED / out_name
    if out.exists():
        return out
    im = Image.open(src).convert('RGB')
    w, h = im.size
    ar = w / h
    if ar > target_ar:                       # largo demais -> corta laterais
        nw = int(h * target_ar)
        im = im.crop(((w - nw) // 2, 0, (w + nw) // 2, h))
    else:                                    # alto demais -> corta topo/base
        nh = int(w / target_ar)
        top = int((h - nh) * 0.45)
        im = im.crop((0, top, w, top + nh))
    im = ImageEnhance.Brightness(im).enhance(factor + 0.25)
    navy = Image.new('RGB', im.size, (29, 45, 84))
    im = Image.blend(im, navy, navy_blend)
    im.save(out)
    return out


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def new(self, footer=True):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        if footer:
            tb = _txbox(s, SW - Inches(3.6), SH - Inches(0.38), Inches(3.3), Inches(0.3))
            p = tb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = f'{self.n}  ·  Defesa de Mestrado — ON'
            _set_font(r, 9, color=GRAY)
        return s

    def title_bar(self, slide, title):
        _rect(slide, 0, 0, SW, TITLE_H, NAVY)
        tb = _txbox(slide, Inches(0.5), Inches(0.08), SW - Inches(1.0), TITLE_H - Inches(0.1))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = tf.paragraphs[0].add_run()
        r.text = title
        _set_font(r, 26, bold=True, color=WHITE)

    def caption(self, slide, text, y=None, size=15, color=DARK, align=PP_ALIGN.CENTER):
        y = y if y is not None else SH - Inches(0.85)
        tb = _txbox(slide, Inches(0.6), y, SW - Inches(1.2), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.alignment = align
        for part in _parse_red(text):
            r = p.add_run()
            r.text = part['t']
            _set_font(r, size, bold=part['red'], color=RED if part['red'] else color)
        return tb


def _parse_red(text):
    """Trechos entre ** ficam vermelhos e bold (números de resultado)."""
    parts, red = [], False
    for chunk in text.split('**'):
        if chunk:
            parts.append({'t': chunk, 'red': red})
        red = not red
    return parts


# ----------------------------------------------------------------------------
# layouts
# ----------------------------------------------------------------------------
def title_slide(d, title, subtitle, author_lines, side_img, logo):
    s = d.new(footer=False)
    _rect(s, 0, 0, SW, SH, NAVY)
    panel_w = SW - Inches(5.0)
    add_pic_fit(s, side_img, panel_w, 0, Inches(5.0), SH)
    tb = _txbox(s, Inches(0.7), Inches(1.5), panel_w - Inches(1.1), Inches(2.6))
    for i, line in enumerate(title.split('\n')):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        r = p.add_run()
        r.text = line
        _set_font(r, 30, bold=True, color=WHITE)
    p = tb.text_frame.add_paragraph()
    p.space_before = Pt(14)
    r = p.add_run()
    r.text = subtitle
    _set_font(r, 18, color=BLUE)
    tb2 = _txbox(s, Inches(0.7), Inches(4.6), panel_w - Inches(1.1), Inches(1.8))
    for i, line in enumerate(author_lines):
        p = tb2.text_frame.paragraphs[0] if i == 0 else tb2.text_frame.add_paragraph()
        r = p.add_run()
        r.text = line
        _set_font(r, 15, color=WHITE, bold=(i == 0))
    add_pic_fit(s, logo, Inches(0.55), SH - Inches(1.5), Inches(1.6), Inches(1.15))
    return s


def section_slide(d, number, title, note=None):
    s = d.new(footer=False)
    _rect(s, 0, 0, SW, SH, NAVY)
    tb = _txbox(s, Inches(1.0), Inches(2.1), Inches(2.6), Inches(2.4))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = number
    _set_font(r, 110, bold=True, color=BLUE)
    tb2 = _txbox(s, Inches(3.6), Inches(2.9), SW - Inches(4.4), Inches(1.6))
    r = tb2.text_frame.paragraphs[0].add_run()
    r.text = title
    _set_font(r, 38, bold=True, color=WHITE)
    if note:
        tb3 = _txbox(s, Inches(3.65), Inches(4.15), SW - Inches(4.6), Inches(0.9))
        r = tb3.text_frame.paragraphs[0].add_run()
        r.text = note
        _set_font(r, 16, italic=True, color=RGBColor(190, 205, 230))
    return s


def image_slide(d, title, img, caption=None):
    s = d.new()
    d.title_bar(s, title)
    bottom = Inches(0.95) if caption else Inches(0.55)
    add_pic_fit(s, img, Inches(0.45), TITLE_H + Inches(0.15),
                SW - Inches(0.9), SH - TITLE_H - bottom - Inches(0.2))
    if caption:
        d.caption(s, caption)
    return s


def two_image_slide(d, title, img_l, img_r, caption=None, weights=(0.5, 0.5)):
    s = d.new()
    d.title_bar(s, title)
    bottom = Inches(0.95) if caption else Inches(0.55)
    area_h = SH - TITLE_H - bottom - Inches(0.2)
    total_w = SW - Inches(1.2)
    wl = int(total_w * weights[0])
    wr = int(total_w * weights[1])
    add_pic_fit(s, img_l, Inches(0.45), TITLE_H + Inches(0.15), wl, area_h)
    add_pic_fit(s, img_r, Inches(0.45) + wl + Inches(0.3), TITLE_H + Inches(0.15), wr, area_h)
    if caption:
        d.caption(s, caption)
    return s


def image_caption_slide(d, title, img, bullets, img_frac=0.62):
    s = d.new()
    d.title_bar(s, title)
    area_h = SH - TITLE_H - Inches(0.7)
    img_w = int((SW - Inches(1.0)) * img_frac)
    add_pic_fit(s, img, Inches(0.45), TITLE_H + Inches(0.15), img_w, area_h)
    tx = Inches(0.65) + img_w
    tb = _txbox(s, tx, TITLE_H + Inches(0.55), SW - tx - Inches(0.5), area_h - Inches(0.8))
    for i, b in enumerate(bullets):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        p.space_after = Pt(14)
        first = True
        for part in _parse_red('•  ' + b):
            r = p.add_run()
            r.text = part['t']
            _set_font(r, 17, bold=part['red'], color=RED if part['red'] else DARK)
            first = False
    return s


def bignum_slide(d, title, numbers, img=None, caption=None):
    """numbers: list[(valor, rotulo, cor)]"""
    s = d.new()
    d.title_bar(s, title)
    n = len(numbers)
    top = TITLE_H + Inches(0.55)
    area_w = SW - Inches(1.2)
    col_w = int(area_w / n)
    for i, (val, label, color) in enumerate(numbers):
        x = Inches(0.6) + i * col_w
        tb = _txbox(s, x, top, col_w, Inches(1.7))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = val
        _set_font(r, 60 if len(val) < 12 else 44, bold=True, color=color)
        tb2 = _txbox(s, x, top + Inches(1.55), col_w, Inches(1.0))
        p2 = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        _set_font(r2, 15, color=DARK)
    if img:
        # termina em 6,55" para nunca invadir a legenda (que começa em 6,65")
        add_pic_fit(s, img, Inches(1.2), top + Inches(2.6), SW - Inches(2.4),
                    SH - top - Inches(3.55))
    if caption:
        d.caption(s, caption)
    return s


def table_slide(d, title, headers, rows, red_cols, caption=None):
    s = d.new()
    d.title_bar(s, title)
    nrows, ncols = len(rows) + 1, len(headers)
    x, y = Inches(0.7), TITLE_H + Inches(0.35)
    w = SW - Inches(1.4)
    h = Inches(0.52) * nrows
    gfx = s.shapes.add_table(nrows, ncols, x, y, w, h)
    table = gfx.table
    widths = [Emu(int(w * f)) for f in (0.10, 0.30, 0.22, 0.19, 0.19)]
    for c, wd in enumerate(widths[:ncols]):
        table.columns[c].width = wd
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = head
        _set_font(r, 15, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        for c, valcell in enumerate(row):
            cell = table.cell(i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 247, 252) if i % 2 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = valcell
            red = c in red_cols
            _set_font(r, 14, bold=red, color=RED if red else DARK)
    if caption:
        d.caption(s, caption)
    return s


def bullets_slide(d, title, bullets, img=None, img_frac=0.34):
    s = d.new()
    d.title_bar(s, title)
    tw = SW - Inches(1.2)
    if img:
        iw = int((SW - Inches(1.0)) * img_frac)
        add_pic_fit(s, img, SW - iw - Inches(0.5), TITLE_H + Inches(0.3),
                    iw, SH - TITLE_H - Inches(1.0))
        tw = SW - iw - Inches(1.5)
    tb = _txbox(s, Inches(0.7), TITLE_H + Inches(0.5), tw, SH - TITLE_H - Inches(1.1))
    for i, b in enumerate(bullets):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        p.space_after = Pt(16)
        for part in _parse_red('•  ' + b):
            r = p.add_run()
            r.text = part['t']
            _set_font(r, 18, bold=part['red'], color=RED if part['red'] else DARK)
    return s


# valores REAIS extraídos de dataset_final.csv (curvas de Umbriel 2020-09-01)
FEATURE_ROWS = [
    ('Umbriel — A. Scheck (positiva)', '0,43', '7,0', '0,142', '−0,74'),
    ('Umbriel — N. Carlson (negativa)', '0,31', '4,4', '0,069', '−0,59'),
]
FEATURE_HEADERS = ['Curva', 'Occ_depth', 'SNR do dip', 'Savgol_std', 'Max_Drawdown']


def arrow_slide(d, title, img_curve, caption, arrow_img, highlight_row):
    """Curva (topo-esq) -> seta vermelha -> tabela NATIVA com os valores reais das
    duas curvas (base-dir); a linha da curva exibida é destacada."""
    s = d.new()
    d.title_bar(s, title)
    top = TITLE_H + Inches(0.12)
    add_pic_fit(s, img_curve, Inches(0.45), top, Inches(7.3), Inches(2.65))
    s.shapes.add_picture(str(arrow_img), Inches(7.75), Inches(2.75),
                         width=Inches(1.4), height=Inches(1.0))
    x, y, w = Inches(4.4), Inches(4.05), Inches(8.4)
    gfx = s.shapes.add_table(3, 5, x, y, w, Inches(1.9))
    table = gfx.table
    table.columns[0].width = Emu(int(w * 0.40))
    for c in range(1, 5):
        table.columns[c].width = Emu(int(w * 0.15))
    for c, head in enumerate(FEATURE_HEADERS):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = head
        _set_font(r, 14, bold=True, color=WHITE)
    for i, row in enumerate(FEATURE_ROWS, start=1):
        hl = (i == highlight_row)
        for c, val in enumerate(row):
            cell = table.cell(i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(252, 233, 236) if hl else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            _set_font(r, 14, bold=hl, color=RED if hl and c else DARK)
    d.caption(s, caption)
    return s


# ----------------------------------------------------------------------------
# o deck
# ----------------------------------------------------------------------------
def build():
    d = Deck()
    F, T, P = A / 'figs', A / 'tese', A / 'part1'

    # A ------------------------------------------------------------ abertura
    title_slide(
        d,
        'Detecção Automatizada de\nOcultações Estelares',
        'Curvas de luz e Machine Learning — Defesa de Mestrado',
        ['Thiago Laidler Vidal Cunha',
         'Orientador: Dr. Júlio Camargo',
         'Observatório Nacional — 2026'],
        side_img=P / 'image12.png', logo=F / 'on_logo.png')

    bullets_slide(d, 'Roteiro', [
        'Ocultações estelares — o fenômeno',
        'O problema e os dados',
        'Machine Learning — sem caixa-preta',
        'Resultados',
        'Quaoar e conclusões'])

    # B --------------------------------------------------- ocultações (fenômeno)
    section_slide(d, '1', 'Ocultações Estelares', 'medir o invisível com uma sombra')

    image_slide(d, 'O que é uma ocultação',
                T / 'curva.png',
                'Um corpo pequeno passa na frente de uma estrela: a luz "pisca". '
                '(Chariklo e seus anéis — ilustração: L. Maquet)')

    two_image_slide(d, 'A sombra varre a Terra',
                    T / 'OccUmbriel_mapa.png', T / 'OccUmbriel_mapa2.png',
                    'Cada observador no caminho da sombra vê o desaparecimento em um instante diferente. (Umbriel, 2020)')

    two_image_slide(d, 'Da câmera à curva de luz',
                    T / 'OccUmbriel_outputExemplo.png', T / 'OccUmbriel_outputExemplo2.png',
                    'Frames de vídeo → fotometria → curva de luz. É este o objeto que analisamos.',
                    weights=(0.42, 0.58))

    two_image_slide(d, 'Por que vale a pena: ciência de precisão',
                    T / 'Cordas atualizadas.png', T / 'Bardecker12.png',
                    'Cada curva vira uma corda; as cordas desenham o corpo — tamanho e forma com precisão de km.')

    # C ------------------------------------------------------ problema e dados
    bignum_slide(d, 'O gargalo é humano',
                 [('milhares', 'de curvas acumuladas em arquivos\ne novas campanhas', BLUE),
                  ('1 par de olhos', 'por curva, na triagem manual', NAVY)],
                 img=F / 'ocultacao_ilustracao.png',
                 caption='Previsões pós-Gaia multiplicaram as campanhas; cidadãos-cientistas observam. A triagem visual não escala.')

    image_slide(d, 'A proposta: uma pipeline de triagem',
                F / 'pipeline.png',
                'Banco de dados → normalização → características (features) → 4 classificadores → probabilidade de ocultação.')

    two_image_slide(d, 'De onde vêm os dados',
                    T / 'fluxo1.png', P / 'image27.jpg',
                    'VizieR (B/occ) + Grupo do Rio → banco SQLite único: todas as curvas e metadados em duas tabelas consultáveis.',
                    weights=(0.38, 0.62))

    bignum_slide(d, 'O conjunto de dados rotulado',
                 [('1693', 'curvas de luz no total', BLUE),
                  ('802', 'positivas (reais, com ocultação)', BLUE),
                  ('891', 'negativas\n702 sintéticas · 186 recortes reais · 3 nativas', BLUE)],
                 caption='Negativas reais quase não são catalogadas — simulador físico (Gomes-Ferrante & Braga-Ribas) preenche a lacuna.')

    arrow_slide(d, 'A curva vira números — exemplo positivo',
                P / 'image25.png',
                'Curva com ocultação → a linha dela na tabela: queda mais funda, SNR maior, dispersão maior.',
                arrow_img=P / 'image21.png', highlight_row=1)

    arrow_slide(d, 'A curva vira números — exemplo negativo',
                P / 'image51.png',
                'A curva sem evento vira números menores nas mesmas colunas — é este contraste que o classificador aprende.',
                arrow_img=P / 'image21.png', highlight_row=2)

    # D ------------------------------------------------------------ ML didático
    section_slide(d, '2', 'Machine Learning — sem caixa-preta',
                  'quatro modelos, todos interpretáveis')

    image_caption_slide(d, 'Classificação supervisionada',
                        T / 'sigmoide.png',
                        ['Aprender com exemplos rotulados (positiva / negativa)',
                         'Saída: probabilidade p de haver ocultação',
                         'Decisão: p ≥ limiar τ (padrão 0,5) → "positiva"'])

    two_image_slide(d, 'Árvores de decisão... e florestas',
                    P / 'image38.png', T / 'RF_img.png',
                    'Uma árvore faz perguntas simples e erra; centenas de árvores votando erram muito menos.')

    image_caption_slide(d, 'Por que comitês funcionam',
                        T / 'manual_2.png',
                        ['Erros independentes se cancelam na votação',
                         '4 modelos: Regressão Logística, Random Forest, XGBoost, CatBoost',
                         'Árvores dão de graça a importância de cada feature'],
                        img_frac=0.52)

    image_caption_slide(d, 'Avaliar sem se enganar',
                        T / 'kfold.png',
                        ['Teste sempre em curvas que o modelo nunca viu',
                         'Validação cruzada (5 dobras) confirma estabilidade',
                         'Métricas: precisão, sensibilidade, F1, AUC-ROC'])

    two_image_slide(d, 'Das 28 às 11 features',
                    P / 'image28.png', F / 'feat_importance_rf.png',
                    'Suavização, profundidade, duração, testes estatísticos... análise de redundância: 28 → 11 sem perder desempenho.')

    # E ------------------------------------------------------------- resultados
    section_slide(d, '3', 'Resultados', 'seis experimentos, um padrão')

    table_slide(d, 'Seis experimentos, desempenho estável',
                ['Exp.', 'Configuração', 'Teste', 'Melhor F1', 'AUC-ROC'],
                [('1', '28 features', 'misto 80/20', '0,9937', '≥ 0,9997'),
                 ('2', '11 features (enxuta)', 'misto 80/20', '0,9906', '≥ 0,9995'),
                 ('3', '11 features', 'só curvas reais', '0,9821', '≥ 0,9950'),
                 ('4', '14 features', 'misto 80/20', '0,9906', '≥ 0,9996'),
                 ('5', '13 features', 'misto 80/20', '0,9905', '≥ 0,9996'),
                 ('6', '12 features (sem K-Means)', 'misto 80/20', '0,9906', '≥ 0,9997')],
                red_cols={3, 4},
                caption='Melhor F1 por experimento: **0,982–0,994** — o desempenho não depende de uma escolha feliz.')

    two_image_slide(d, 'Separação das classes (Experimento 1)',
                    F / 'exp1_roc.png', F / 'exp1_confusion.png',
                    'Curvas ROC coladas no canto ideal; erros contados nos dedos. Os 4 modelos, estatisticamente equivalentes (McNemar).')

    two_image_slide(d, 'E se o teste tiver só curvas reais?',
                    F / 'exp3real_roc.png', F / 'exp3real_confusion.png',
                    'Teste 100% real: F1 ≈ **0,98** (queda de ~1 ponto) — consistente com aprender a física, não o simulador.')

    image_caption_slide(d, 'Importância não é insubstituibilidade',
                        T / 'histograma_kmeans_por_classe.png',
                        ['K-Means: a feature "mais importante" das árvores',
                         'Removida → F1 não cai (features colineares compensam)',
                         'Lição: importância alta ≠ informação exclusiva'],
                        img_frac=0.58)

    two_image_slide(d, 'Onde o modelo erra',
                    T / 'RF_curva_FN_1.png', T / 'RF_curva_FN_2.png',
                    'Falsos negativos são curvas que até olhos humanos discutiriam — ruído alto, quedas rasas.')

    image_caption_slide(d, 'O limiar τ: sensibilidade quase de graça',
                        F / 'metrics_vs_threshold.png',
                        ['Perder um evento custa MUITO mais que revisar um alarme',
                         'Baixar τ: sensibilidade de **99,4%** no teste (+5 alarmes em 179 negativas)',
                         'Pós-processamento: sem retreinar nada'],
                        img_frac=0.52)

    # F ------------------------------------------------------------------ Quaoar
    section_slide(d, '4', 'Prova de fogo: Quaoar',
                  'curva real, externa ao treino — cedida por Pereira et al. (2023)')

    image_slide(d, 'A curva completa',
                T / 'quaoar_test.png',
                'Corpo principal + anéis Q1R e Q2R (o fino). O modelo nunca viu Quaoar: p ≈ **0,99** para a curva inteira.')

    image_slide(d, 'Caçando estruturas por recortes',
                T / 'Quaoar_XgBoost_full_pict.png',
                'Janela a janela: anéis fortes passam (p = **0,96–0,999**)... mas as travessias do anel fino Q2R "reprovam" (p = **0,04–0,08**).')

    bignum_slide(d, 'O sinal fraco ainda é sinal',
                 [('86×', 'anel fino real vs. ruído parecido\n(0,043 vs 0,0005)', RED),
                  ('τ = 0,03', 'recupera as duas travessias\ndo anel fino Q2R', NAVY),
                  ('0', 'falsos alarmes nesta curva', RED)],
                 img=T / 'quaoar_recorte_tau_ajustado.png',
                 caption='O que importa não é a probabilidade absoluta — é a separação entre evento e ruído.')

    image_slide(d, 'Lição estrutural: janelas maiores diluem',
                T / 'quaoar_recortes_janelas_expandidas.png',
                'Ampliar a janela REDUZ a probabilidade: features globais diluem quedas curtas → motivação para janela deslizante.')

    # G --------------------------------------------------------------- fechamento
    bullets_slide(d, 'Dificuldades e como foram enfrentadas', [
        'Dados heterogêneos e dispersos → banco SQLite unificado',
        'Só **3** negativas reais nativas → simulador físico + recortes de curvas reais',
        'Redundância de features → análise sistemática: 28 → 11',
        'Eventos curtos e rasos → limiar τ ajustável por campanha'],
        img=T / 'fluxo2.png')

    s = bignum_slide(d, 'Conclusões',
                     [('0,98–0,99', 'F1 robusto em 6 experimentos\n(4 modelos, teste real incluído)', RED),
                      ('99,4%', 'sensibilidade com τ ajustado\nsem retreinamento', RED),
                      ('86×', 'separação anel real vs. ruído\nem dado externo (Quaoar)', RED)],
                     caption='Uma ferramenta reproduzível de triagem e priorização — pronta para revisitar arquivos de curvas.')

    bullets_slide(d, 'Trabalhos futuros', [
        'Janela deslizante + features locais (detectar anéis finos direto)',
        'Ampliar negativas reais nativas (colaborações e campanhas)',
        'Redes neurais (CNN / ODNet) sobre a série bruta',
        'Triagem em tempo quase-real nas campanhas do Grupo do Rio'],
        img=F / 'pipeline.png')

    s = d.new(footer=False)
    _rect(s, 0, 0, SW, SH, NAVY)
    tb = _txbox(s, Inches(1.0), Inches(2.2), SW - Inches(2.0), Inches(1.4))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'Obrigado.'
    _set_font(r, 54, bold=True, color=WHITE)
    tb2 = _txbox(s, Inches(1.0), Inches(3.7), SW - Inches(2.0), Inches(2.2))
    linhas = ['Dr. Júlio Camargo (orientação) · Grupo do Rio',
              'W. Gomes-Ferrante & F. Braga-Ribas (simulador) · C. L. Pereira (dados de Quaoar)',
              'Banca examinadora · família e amigos']
    for i, line in enumerate(linhas):
        p = tb2.text_frame.paragraphs[0] if i == 0 else tb2.text_frame.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = line
        _set_font(r, 16, color=RGBColor(190, 205, 230))
    add_pic_fit(s, F / 'on_logo.png', SW - Inches(2.3), SH - Inches(1.9), Inches(1.7), Inches(1.4))

    # ------------------------------------------------------------------ backups
    section_slide(d, '+', 'Material de apoio', 'slides extras para perguntas')

    table_slide(d, 'Análise complementar: e sem ML, dava?',
                ['', 'Classificador', 'Teste', 'F1', 'AUC'],
                [('—', 'Limiar em 1 feature (Occ_depth)', 'misto 80/20', '0,905', '0,972'),
                 ('—', 'Limiar em 1 feature (Occ_depth)', 'só reais', '0,948', '0,966'),
                 ('ML', 'XGBoost (11 features)', 'misto 80/20', '0,9906', '0,9995'),
                 ('ML', 'XGBoost (11 features)', 'só reais', '0,9821', '0,9968')],
                red_cols={3, 4},
                caption='Uma régua já vai longe (a física é forte) — o ML corta o erro restante em **~10×**. [análise complementar, pós-dissertação]')

    image_slide(d, 'Teste cego: Umbriel (seminário, 2025)',
                P / 'image51.png',
                'A única divergência do teste cego (18 curvas: 15 pos / 3 neg) foi esta curva, de rótulo discutível — era mesmo negativa?')

    two_image_slide(d, 'Recortes de Quaoar em detalhe',
                    T / 'Quaoar_XgBoost_2st3st_cut.png', T / 'Quaoar_XgBoost_5st_cut.png',
                    'Anel fino real (esq.) vs. anel denso (dir.) — probabilidades do XGBoost janela a janela.')

    image_slide(d, 'Física do perfil: modelo de curva',
                T / 'ModelosSORA.png',
                'Difração de Fresnel + diâmetro estelar + tempo de integração — o modelo físico por trás da forma da queda (SORA).')

    d.prs.save(OUT)
    print(f'OK: {OUT.name} gerado com {d.n} slides.')
    return d


if __name__ == '__main__':
    build()
